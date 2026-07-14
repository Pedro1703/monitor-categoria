#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Verificador de diseño del PPT. Corre SIEMPRE después de generar el informe.

Un PPT que "compila" no es un PPT que se ve bien. Estos son los errores que ya pasaron
y que este script ahora detecta antes de que lleguen a una reunión:

  1. TEXTO QUE SE SUPERPONE con otra forma (el caso de las etiquetas largas metiéndose
     dentro de las barras del ranking).
  2. IMÁGENES que invaden la zona del pie y tapan la nota metodológica.
  3. ELEMENTOS FUERA DEL BORDE de la slide.
  4. TEXTO CORTADO EN VERTICAL: un textbox más angosto que su contenido, que PowerPoint
     parte carácter por carácter (el "01" que salía 0 arriba y 1 abajo).
  5. Slides con datos SIN nota metodológica.

Si algo falla, sale con código 1 y el informe NO se da por bueno.

    python3 verificar_ppt.py
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Emu

HERE = os.path.dirname(os.path.abspath(__file__))
PPT = os.path.join(HERE, "Informe_Ciudadana.pptx")

# Ancho medio de un carácter, en pulgadas por punto de tamaño de fuente.
# Sirve para estimar si un texto entra en su caja.
CHAR_W = 0.0075
ZONA_PIE = Inches(6.62)      # de acá para abajo va la nota metodológica


def rects_chocan(a, b, holgura=Emu(20000)):
    """¿Se superponen dos rectángulos? Con una holgura para no marcar roces mínimos."""
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    return not (ax2 - holgura <= bx1 or bx2 - holgura <= ax1 or
                ay2 - holgura <= by1 or by2 - holgura <= ay1)


def ancho_texto(run):
    """Ancho estimado del texto de un run, en EMU."""
    size = run.font.size.pt if run.font.size else 12
    return Inches(len(run.text) * CHAR_W * size)


def verificar(ruta=PPT):
    if not os.path.exists(ruta):
        print("ERROR: no existe %s" % ruta, file=sys.stderr)
        return 1

    prs = Presentation(ruta)
    W, H = prs.slide_width, prs.slide_height
    fallos = []

    for i, s in enumerate(prs.slides, 1):
        textos, formas, imagenes = [], [], []
        tiene_nota = False
        es_datos = False

        for sh in s.shapes:
            if sh.left is None or sh.top is None:
                continue
            caja = (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)

            # ── fuera de borde
            if (sh.left < -Emu(50000) or sh.top < -Emu(50000)
                    or sh.left + sh.width > W + Emu(120000)
                    or sh.top + sh.height > H + Emu(120000)):
                fallos.append("slide %d · elemento fuera del borde" % i)

            if sh.shape_type == 13:                      # imagen
                imagenes.append((sh, caja))
                if sh.top + sh.height > ZONA_PIE and sh.height > Inches(1):
                    fallos.append("slide %d · una imagen invade la zona del pie "
                                  "(tapa la nota metodológica)" % i)
                continue

            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text
                if "© 2026" in t:
                    es_datos = True
                if t.strip().startswith("Metodología ·"):
                    tiene_nota = True
                    continue                              # el pie puede rozar la línea

                # ── texto cortado en vertical
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if not r.text.strip():
                            continue
                        need = ancho_texto(r)
                        if sh.text_frame.word_wrap and need > sh.width * 2.2 \
                                and len(r.text) <= 6:
                            fallos.append("slide %d · «%s» puede partirse en vertical "
                                          "(caja más angosta que el texto)" % (i, r.text))
                        # ── texto que se desborda de su caja hacia la derecha
                        if not sh.text_frame.word_wrap and need > sh.width * 1.15:
                            textos.append((sh, (sh.left, sh.top,
                                                sh.left + need, sh.top + sh.height), r.text))
                            continue
                textos.append((sh, caja, sh.text_frame.text[:28]))
            elif sh.shape_type is not None:
                formas.append((sh, caja))

        # ── texto encima de una forma (barras del ranking, tarjetas)
        for _, caja_t, txt in textos:
            for _, caja_f in formas:
                if rects_chocan(caja_t, caja_f):
                    fallos.append("slide %d · el texto «%s» se superpone con una forma"
                                  % (i, txt.strip()[:30]))
                    break

        # ── texto encima de una imagen
        for _, caja_t, txt in textos:
            for _, caja_i in imagenes:
                if rects_chocan(caja_t, caja_i):
                    fallos.append("slide %d · el texto «%s» se superpone con una imagen"
                                  % (i, txt.strip()[:30]))
                    break

        if es_datos and not tiene_nota:
            fallos.append("slide %d · tiene datos pero NO tiene nota metodológica" % i)

    # Deduplicar sin perder el orden
    vistos, unicos = set(), []
    for f in fallos:
        if f not in vistos:
            vistos.add(f)
            unicos.append(f)

    print("Verificando %s (%d slides)…\n" % (os.path.basename(ruta), len(prs.slides._sldIdLst)))
    if unicos:
        for f in unicos:
            print("  ✗ %s" % f)
        print("\n%d problema(s) de diseño. El informe NO está listo." % len(unicos))
        return 1
    print("  ✓ Sin superposiciones")
    print("  ✓ Sin elementos fuera de borde")
    print("  ✓ Sin texto partido en vertical")
    print("  ✓ Todas las slides con datos tienen nota metodológica")
    print("\nEl informe está limpio.")
    return 0


if __name__ == "__main__":
    sys.exit(verificar())
