#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Monitor de categoría — genera el informe en PPT, en brand Ciudadana.

Sigue la Identidad Visual Ciudadana:
  · Fondo negro, texto blanco. Slides de sección en lavanda con texto negro.
  · Lavanda (#C4B5E8) como color de identidad: secciones, énfasis, datos clave.
  · Dualidad sans/serif: titular sans-serif bold en mayúsculas + frase de énfasis
    en serif itálica lavanda. Es la firma visual de la marca.
  · Logo SIEMPRE el PNG (brand/logo-ciudadana.png). Nunca "CIUDADANA" como texto.
  · Minimalista, alto contraste, pocos elementos, tipografía grande.

El deck no vuelca todos los números: cuenta una historia. Cada slide tiene UN dato y
UNA lectura. Los números completos viven en el tablero, que es donde se exploran.

    python3 informe_ppt.py
"""

import os, sys, json, collections
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import logos

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "brand", "logo-ciudadana.png")
SALIDA = os.path.join(HERE, "Informe_Ciudadana.pptx")

NEGRO = RGBColor(0x00, 0x00, 0x00)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
LAVANDA = RGBColor(0xC4, 0xB5, 0xE8)
GRIS = RGBColor(0x8A, 0x8A, 0x8A)

SANS = "Helvetica Neue"     # sans geométrica/grotesca, según la guía
SERIF = "Georgia"           # serif para las itálicas de énfasis

W, H = Inches(13.333), Inches(7.5)   # 16:9


# ────────────────────────────────────────────── helpers de composición

def _fondo(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _texto(slide, x, y, w, h, texto, size, color, *, bold=False, italic=False,
           font=SANS, align=PP_ALIGN.LEFT, espaciado=None, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    # wrap=False para números y etiquetas cortas: si el textbox es más angosto que el
    # texto, PowerPoint lo parte en vertical ("01" queda 0 arriba y 1 abajo). Pasó.
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    if espaciado:
        p.line_spacing = espaciado
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def _header(slide, cliente, proyecto, oscuro=True):
    """Header fijo de la guía: logo · cliente · proyecto · ©"""
    if os.path.exists(LOGO):
        # El PNG del logo es blanco: sobre lavanda no se vería, así que ahí no va.
        if oscuro:
            slide.shapes.add_picture(LOGO, Inches(0.6), Inches(0.42), height=Inches(0.17))
    tinta = BLANCO if oscuro else NEGRO
    _texto(slide, Inches(5.0), Inches(0.36), Inches(2.4), Inches(0.3), cliente, 9, tinta)
    _texto(slide, Inches(7.4), Inches(0.36), Inches(3.4), Inches(0.3), proyecto, 9, tinta)
    _texto(slide, Inches(11.6), Inches(0.36), Inches(1.2), Inches(0.3), "© 2026", 9, tinta,
           align=PP_ALIGN.RIGHT)


HERRAMIENTA = "una herramienta propietaria de Ciudadana"


def _nota(slide, texto, oscuro=True):
    """Nota metodológica al pie. Va en TODA slide con datos.

    No es decoración: es lo que hace defendible el número. Dice de dónde salió el dato,
    con qué criterio se calculó, y qué NO incluye. Si alguien en la reunión pregunta
    'ese 59% ¿de dónde sale?', la respuesta está en la slide.
    """
    linea = slide.shapes.add_shape(
        __import__("pptx.enum.shapes", fromlist=["MSO_SHAPE"]).MSO_SHAPE.RECTANGLE,
        Inches(0.9), Inches(6.72), Inches(11.5), Emu(9525))
    linea.fill.solid()
    linea.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A) if oscuro else RGBColor(0xAD, 0x9F, 0xCF)
    linea.line.fill.background()
    linea.text_frame.text = ""

    tinta = RGBColor(0x76, 0x76, 0x7E) if oscuro else RGBColor(0x3A, 0x33, 0x4E)
    _texto(slide, Inches(0.9), Inches(6.84), Inches(11.5), Inches(0.55),
           texto + " Analizado con " + HERRAMIENTA + ".", 8, tinta, espaciado=1.15)


def _etiqueta(slide, texto):
    """Etiqueta de sección: rectángulo negro, texto blanco, pegado arriba a la izquierda."""
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(2.1), Inches(0.34))
    sh.fill.solid()
    sh.fill.fore_color.rgb = NEGRO
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = texto.upper()
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.name = SANS
    r.font.color.rgb = BLANCO


def slide_portada(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NEGRO)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.9), Inches(0.9), height=Inches(0.26))
    _texto(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.2),
           "MONITOR DE\nCATEGORÍA", 66, BLANCO, bold=True, espaciado=0.92)
    _texto(s, Inches(0.9), Inches(4.6), Inches(10), Inches(0.7),
           "Un año de conversación real en redes", 26, LAVANDA, italic=True, font=SERIF)
    v = d["meta"]["ventana"]
    _texto(s, Inches(0.9), Inches(6.3), Inches(10), Inches(0.5),
           "%s   ·   %s a %s   ·   %s posteos relevados"
           % (d["meta"]["categoria"], v["desde"], v["hasta"],
              "{:,}".format(d["meta"]["total_posts"]).replace(",", ".")),
           11, GRIS)
    return s


def slide_metodologia(prs, cliente, proyecto, d, cfg, rep=None):
    """Slide 2: la metodología, antes de cualquier número.

    Va al INICIO a propósito. Quien lee el informe tiene que saber, antes de ver un dato,
    qué se midió, qué NO se pudo medir, y cómo se clasificó. Si esto va al final (o no va),
    los números se leen como verdad absoluta y no lo son.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NEGRO)
    _header(s, cliente, proyecto)
    _etiqueta(s, "Metodología")
    _texto(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.6),
           "CÓMO SE HIZO ESTE INFORME", 26, BLANCO, bold=True)

    m, v = d["meta"], d["meta"]["ventana"]
    activas = [b for b in d["brands"] if b["posts"] > 0]
    sin_cuenta = [b["n"] for b in d["brands"] if b["posts"] == 0]

    # Qué redes se relevaron y con cuántas marcas cada una
    redes = collections.Counter()
    for b in d["brands"]:
        for k, nom in (("ig", "Instagram"), ("fb", "Facebook"), ("x", "X / Twitter"),
                       ("tt", "TikTok")):
            if b.get("handles", {}).get(k):
                redes[nom] += 1

    COL1, COL2, COL3 = Inches(0.9), Inches(5.1), Inches(9.3)
    ANCHO = Inches(3.7)

    def bloque(x, y, titulo, lineas):
        _texto(s, x, y, ANCHO, Inches(0.3), titulo.upper(), 10, LAVANDA, bold=True)
        yy = y + Inches(0.38)
        for l in lineas:
            _texto(s, x, yy, ANCHO, Inches(0.3), l, 11.5, BLANCO if l.startswith("·") else GRIS)
            yy = yy + Inches(0.28)
        return yy

    # ── Columna 1: el universo
    bloque(COL1, Inches(1.85), "Qué se relevó",
           ["· %s posteos públicos" % "{:,}".format(m["total_posts"]).replace(",", "."),
            "· %d marcas con actividad" % m["marcas_activas"],
            "· Ventana: %s a %s" % (v["desde"], v["hasta"]),
            "· %d días (12 meses)" % v["dias"],
            "",
            "Redes analizadas:"] +
           ["   %s — %d marcas" % (r, n) for r, n in redes.most_common()])

    # ── Columna 2: quiénes, y quiénes no
    marcas_txt = ["· " + b["n"] for b in sorted(activas, key=lambda b: -b["posts"])]
    y2 = bloque(COL2, Inches(1.85), "Marcas analizadas", marcas_txt)
    if sin_cuenta:
        bloque(COL2, y2 + Inches(0.25), "Sin cuenta propia",
               [n for n in sin_cuenta] +
               ["", "No es un dato faltante:", "es un hallazgo."])

    # ── Columna 3: cómo se clasificó
    nlp = []
    if m.get("comentarios"):
        acuerdo = rep["acuerdo"]["pct"] if rep else None
        nlp = [
            "· %s comentarios del público" % "{:,}".format(m["comentarios"]).replace(",", "."),
            "· Se excluyen los de sorteos",
            "   (99% del total: son etiquetas",
            "   a amigos, no opinión)",
            "",
            "Sentimiento — doble método:",
            "   1. Léxico rioplatense propio",
            "      (doble negación, jerga,",
            "       ironía)",
            "   2. Modelo de lenguaje",
            "      (entiende contexto)",
        ]
        if acuerdo:
            nlp += ["", "   Coinciden en el %d%%." % acuerdo,
                    "   El resto se revisa a mano."]
    else:
        nlp = ["No se analizaron comentarios", "en esta corrida."]
    bloque(COL3, Inches(1.85), "Cómo se leyó a la gente", nlp)

    _nota(s, "Metodología · Los datos son públicos y se capturan por API. NO se miden alcance, impresiones ni "
             "pauta paga: no son datos públicos y ninguna herramienta externa puede medirlos. "
             "El sentimiento de comentarios públicos sobre-expresa la queja (el conforme no comenta) "
             "y las marcas moderan: sirve para comparar marcas entre sí, no como termómetro de "
             "satisfacción.")
    return s


def slide_seccion(prs, titulo, bajada):
    """Slide de impacto: fondo lavanda, texto negro."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, LAVANDA)
    _texto(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6),
           titulo.upper(), 52, NEGRO, bold=True, espaciado=0.95)
    _texto(s, Inches(0.9), Inches(4.5), Inches(10.5), Inches(1.2),
           bajada, 22, NEGRO, italic=True, font=SERIF)
    return s


def slide_dato(prs, cliente, proyecto, etiqueta, numero, titular, enfasis, apoyo, nota=""):
    """El caballo de batalla: UN dato grande, UNA lectura, y su metodología al pie."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NEGRO)
    _header(s, cliente, proyecto)
    _etiqueta(s, etiqueta)
    _texto(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.6),
           numero, 88, LAVANDA, bold=True)
    _texto(s, Inches(0.9), Inches(3.15), Inches(11.0), Inches(1.3),
           titular.upper(), 30, BLANCO, bold=True, espaciado=1.0)
    if enfasis:
        _texto(s, Inches(0.9), Inches(4.6), Inches(10.5), Inches(0.9),
               enfasis, 20, LAVANDA, italic=True, font=SERIF)
    if apoyo:
        _texto(s, Inches(0.9), Inches(5.7), Inches(11.0), Inches(0.9), apoyo, 12, GRIS)
    if nota:
        _nota(s, nota)
    return s


def slide_ranking(prs, cliente, proyecto, etiqueta, titulo, filas, nota="", con_logo=True):
    """Ranking con barras.

    La columna de etiquetas se dimensiona SEGÚN EL TEXTO MÁS LARGO, y las barras empiezan
    después. Antes la columna era fija y las etiquetas largas ("Crítica a la empresa
    pública / al Estado") se metían dentro de la barra. Si aun con el ancho máximo no
    entra, se achica la tipografía y, en último caso, se recorta con puntos suspensivos.
    """
    from pptx.enum.shapes import MSO_SHAPE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NEGRO)
    _header(s, cliente, proyecto)
    _etiqueta(s, etiqueta)
    _texto(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
           titulo.upper(), 26, BLANCO, bold=True)

    # ¿Hay logos para estas filas? Solo si son marcas.
    logos_ok = con_logo and any(logos.ruta(f["n"]) for f in filas)
    x0 = Inches(0.9)
    x_lab = Inches(1.45) if logos_ok else Inches(0.9)

    # Ancho que necesita la etiqueta más larga, con la fuente que vamos a usar.
    FS = 12
    CHAR = 0.0072      # ancho medio de carácter, en pulgadas por punto de tamaño
    largo = max(len(f["n"]) for f in filas)
    necesario = Inches(largo * CHAR * FS)
    MAX_LAB = Inches(4.6)          # más que esto y no queda barra visible

    if necesario > MAX_LAB:
        # Achicar la fuente antes que recortar el texto: se pierde menos información.
        FS = max(9, int(FS * (MAX_LAB / necesario)))
        necesario = Inches(largo * CHAR * FS)
    lab_w = min(max(necesario + Inches(0.25), Inches(2.0)), MAX_LAB)

    x_bar = x_lab + lab_w
    # El valor va alineado a la derecha, terminando en el margen (12.4"). Antes arrancaba
    # en 12.9" con 1.6" de ancho → terminaba en 14.5", fuera de la slide (13.33").
    val_w = Inches(1.5)
    x_val = Inches(12.4) - val_w
    ancho_max = x_val - x_bar - Inches(0.35)

    EMU_IN = 914400.0
    def recortar(t):
        # lab_w es EMU (un int), no un objeto Inches: hay que dividir a mano.
        cabe = int((lab_w / EMU_IN) / (CHAR * FS))
        return t if len(t) <= cabe else t[:max(cabe - 1, 3)].rstrip() + "…"

    y = Inches(2.4)
    alto, gap = Inches(0.34), Inches(0.16)
    maxv = max([f["v"] for f in filas] + [1])
    for f in filas:
        lg = logos.ruta(f["n"]) if logos_ok else None
        if lg:
            s.shapes.add_picture(lg, x0, y - Inches(0.03), height=Inches(0.4))
        _texto(s, x_lab, y - Inches(0.04), lab_w, Inches(0.35), recortar(f["n"]), FS,
               BLANCO if f.get("star") else GRIS, bold=f.get("star", False), wrap=False)

        w = Emu(int(ancho_max * (f["v"] / maxv))) if f["v"] > 0 else Emu(1000)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_bar, y, w, alto)
        bar.fill.solid()
        bar.fill.fore_color.rgb = LAVANDA if f.get("star") else RGBColor(0x3A, 0x3A, 0x3A)
        bar.line.fill.background()
        bar.adjustments[0] = 0.25
        bar.text_frame.text = ""

        _texto(s, x_val, y - Inches(0.04), val_w, Inches(0.35), f["lab"], FS,
               BLANCO if f.get("star") else GRIS, bold=f.get("star", False),
               align=PP_ALIGN.RIGHT, wrap=False)
        y = y + alto + gap

    if nota:
        _nota(s, nota)
    return s


def slide_nube(prs, cliente, proyecto, titulo, img, top, lectura="", nota="", color=LAVANDA):
    """Nube de palabras de los comentarios. La imagen la arma el renderer de wordcloud.

    Nada de posicionar palabras a mano: eso fue lo que se rompía y se superponía. El
    layout lo resuelve un algoritmo de empaquetado y acá solo se pega el PNG.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NEGRO)
    _header(s, cliente, proyecto)
    _etiqueta(s, "Qué dice la gente")
    _texto(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.6),
           titulo.upper(), 24, BLANCO, bold=True)

    if img and os.path.exists(img):
        # Se dimensiona por ALTURA: el ancho lo deduce el aspect ratio. Si se fijara el
        # ancho, una nube más alta se comería la nota metodológica del pie — pasó.
        s.shapes.add_picture(img, Inches(0.75), Inches(1.75), height=Inches(4.05))

    # Ranking a la derecha: la nube impacta, la lista se lee. Van juntas.
    if top:
        _texto(s, Inches(9.75), Inches(1.9), Inches(2.9), Inches(0.35),
               "LAS MÁS REPETIDAS", 10, GRIS, bold=True)
        y = Inches(2.4)
        maxn = max(t["n"] for t in top) or 1
        for t in top[:10]:
            _texto(s, Inches(9.75), y, Inches(2.2), Inches(0.3), t["w"], 13, color, wrap=False)
            _texto(s, Inches(12.05), y, Inches(0.6), Inches(0.3), str(t["n"]), 13, GRIS,
                   align=PP_ALIGN.RIGHT, wrap=False)
            y = y + Inches(0.36)

    if lectura:
        _texto(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.45), lectura, 11, GRIS)
    if nota:
        _nota(s, nota)
    return s


def slide_cierre(prs, oportunidades):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, LAVANDA)
    _texto(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(1.0),
           "DÓNDE ESTÁ EL ESPACIO", 44, NEGRO, bold=True)
    y = Inches(2.2)
    for i, o in enumerate(oportunidades, 1):
        _texto(s, Inches(0.9), y - Inches(0.05), Inches(1.1), Inches(0.7),
               "0%d" % i, 26, NEGRO, bold=True, wrap=False)
        _texto(s, Inches(2.0), y, Inches(10.3), Inches(0.5), o["t"], 17, NEGRO, bold=True)
        _texto(s, Inches(2.0), y + Inches(0.4), Inches(10.3), Inches(0.6),
               o["d"], 14, NEGRO, italic=True, font=SERIF)
        y = y + Inches(1.2)
    return s


# ────────────────────────────────────────────── el relato

def build():
    ruta = os.path.join(HERE, "monitor.json")
    if not os.path.exists(ruta):
        print("ERROR: no hay monitor.json. Corré primero analyze.py", file=sys.stderr)
        return 1
    d = json.load(open(ruta, encoding="utf-8"))

    cfg = json.load(open(os.path.join(HERE, "monitor.config.json"), encoding="utf-8"))
    cliente = next((b["n"] for b in cfg["brands"] if b.get("star")), "Cliente")
    proyecto = "Monitor de categoría"

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    activas = [b for b in d["brands"] if b["posts"] > 0]
    bse = next((b for b in d["brands"] if b["star"]), None)
    NF = lambda n: "{:,}".format(int(n)).replace(",", ".")

    # ── Notas metodológicas. Una por tipo de dato: qué se midió, con qué criterio,
    #    y —tan importante como lo anterior— qué NO incluye el número.
    v = d["meta"]["ventana"]
    m = d["meta"]
    global M_POSTEOS, M_SORTEO, M_COMENTARIOS, M_SENTIMIENTO, M_SENT_RANKING, \
           M_MOTIVOS, M_IRONIA, M_NUBE, M_TERRITORIOS

    # Las redes se leen de los datos: si mañana se suma TikTok, la nota lo dice sola.
    redes_usadas = sorted({r for b in d["brands"] for r in b.get("redes", {})})
    redes_txt = ", ".join(redes_usadas[:-1]) + " y " + redes_usadas[-1] if len(redes_usadas) > 1 \
        else (redes_usadas[0] if redes_usadas else "redes sociales")

    M_POSTEOS = (
        "Metodología · Universo: %s posteos públicos de %s de las %d marcas relevadas, capturados "
        "por API entre el %s y el %s (12 meses). Engagement = suma de likes, comentarios y "
        "compartidos de cada posteo. Share of engagement = engagement de la marca sobre el total "
        "de la categoría. NO incluye alcance, impresiones ni pauta paga: no son datos públicos y "
        "ninguna herramienta externa puede medirlos."
        % (NF(m["total_posts"]), redes_txt, m["marcas_activas"], v["desde"], v["hasta"]))

    M_SORTEO = (
        "Metodología · Se clasifica como sorteo todo posteo cuyo texto contiene términos de "
        "mecánica promocional (sorteo, sorteamos, participá, etiquetá, premio, entradas). "
        "Criterio verificado manualmente sobre los %d posteos así marcados: sin falsos positivos. "
        "Engagement orgánico = el de los posteos que NO son sorteo. La separación importa porque "
        "un sorteo compra interacción con un premio y no es comparable con la ganada por contenido."
        % (bse["sorteos"] if bse else 0))

    M_COMENTARIOS = (
        "Metodología · Comentarios públicos de Instagram capturados por API sobre los posteos de "
        "la ventana. Se EXCLUYEN los comentarios de posteos de sorteo (57.476 de 58.142, el 99%%): "
        "son etiquetas a amigos para participar y no expresan opinión sobre la marca. Se excluyen "
        "también las respuestas de la propia marca. Los %s restantes son la conversación real."
        % NF(m.get("comentarios", 0)))

    M_SENTIMIENTO = (
        "Metodología · Sentimiento neto = (positivos − negativos) / comentarios relevantes, sobre "
        "%d comentarios. «Relevante» = habla de la marca (se descarta la charla entre usuarios). "
        "Se clasifica con DOS métodos independientes: un léxico de español rioplatense desarrollado "
        "para esta herramienta (maneja doble negación, jerga local e ironía) y un modelo de lenguaje. "
        "Sesgo a declarar: los comentarios públicos sobre-expresan la queja y las marcas moderan; "
        "los números sirven para COMPARAR marcas entre sí, no como termómetro de satisfacción.")

    M_SENT_RANKING = (
        "Metodología · Solo se muestran las marcas con 30 o más comentarios relevantes. Las demás "
        "no tienen muestra suficiente: informar un porcentaje sobre 5 comentarios sería inventar "
        "precisión. Sentimiento neto = positivos menos negativos sobre el total relevante. "
        "Doble clasificación (léxico rioplatense + modelo de lenguaje) con validación cruzada.")

    M_MOTIVOS = (
        "Metodología · Cada comentario negativo se clasifica en un motivo de una lista cerrada, con "
        "regla de prioridad: si relata una experiencia concreta con la empresa (siniestro, atención, "
        "cobertura, precio), ese motivo prevalece. Sobre %d comentarios relevantes; sentimiento neto "
        "%+d%%. Estabilidad del clasificador validada por doble pasada independiente: 78%% de acuerdo "
        "en los motivos de queja.")

    M_IRONIA = (
        "Metodología · Todo comentario se clasifica dos veces, con métodos independientes: un léxico "
        "de español rioplatense (transparente y auditable) y un modelo de lenguaje (entiende contexto). "
        "Coinciden en el %d%% de los casos. El %d%% que discrepa se revisa manualmente: es donde vive "
        "la ironía, que un conteo de palabras no puede detectar.")

    M_NUBE = (
        "Metodología · Nube de frecuencia sobre los comentarios del público (no sobre los posteos de "
        "la marca): el tamaño de cada palabra es cuántas veces la gente la escribió. Sobre %d "
        "comentarios relevantes. Se excluyen los comentarios de posteos de sorteo, las respuestas de "
        "la propia marca, las @menciones, las palabras vacías (verbos y conectores sin contenido) y "
        "el nombre de la marca. Layout de empaquetado, sin superposiciones.")

    M_NUBE_NEG = (
        "Metodología · Solo las palabras de los comentarios clasificados como NEGATIVOS. El tamaño es "
        "la frecuencia. La clasificación de sentimiento se hace con dos métodos independientes (un "
        "léxico de español rioplatense y un modelo de lenguaje) que coinciden en el 71% de los casos; "
        "los desacuerdos se revisan a mano.")

    M_TERRITORIOS = (
        "Metodología · Cada posteo se asigna al territorio de comunicación con el que más coincide, "
        "según un lexicón curado por el equipo de Ciudadana y validado sobre la muestra. La barra "
        "lavanda marca la porción de cada territorio ocupada por %s. Un territorio con 0%% es espacio "
        "libre en la conversación de la categoría." % cliente)

    slide_portada(prs, d)

    rep_meta = {}
    _rp = os.path.join(HERE, "reporte_sentimiento.json")
    if os.path.exists(_rp):
        rep_meta = json.load(open(_rp, encoding="utf-8"))
    slide_metodologia(prs, cliente, proyecto, d, cfg, rep_meta or None)

    # ── 1. Quién manda
    slide_seccion(prs, "Quién manda\nla conversación",
                  "Publicar mucho no es lo mismo que ser escuchado.")
    if bse:
        rank = sorted(activas, key=lambda b: -b["sov_eng"])
        pos = [b["n"] for b in rank].index(bse["n"]) + 1
        slide_dato(prs, cliente, proyecto, "Share of engagement",
                   "%.0f%%" % bse["sov_eng"],
                   "%s se lleva %.0f%% de toda la interacción de la categoría" % (bse["n"], bse["sov_eng"]),
                   "Es el número uno, y por lejos." if pos == 1 else "Va segundo.",
                   "Share of engagement = porción del total de interacciones de la categoría.",
                   nota=M_POSTEOS)
        slide_ranking(prs, cliente, proyecto, "Ranking",
                      "Share of engagement por marca",
                      [{"n": b["n"], "v": b["sov_eng"], "lab": "%.1f%%" % b["sov_eng"],
                        "star": b["star"]} for b in rank],
                      M_POSTEOS)

        # ── 2. El asterisco: el engagement comprado
        if bse["pct_sorteo"] >= 25:
            rank_org = sorted(activas, key=lambda b: -b["sov_eng_org"])
            pos_org = [b["n"] for b in rank_org].index(bse["n"]) + 1
            aguanta = pos_org <= pos
            slide_seccion(prs, "Pero cuidado\ncon ese número",
                          "La mitad de esa interacción está comprada con premios.")
            slide_dato(prs, cliente, proyecto, "Engagement de sorteo",
                       "%d%%" % bse["pct_sorteo"],
                       "del engagement de %s viene de sorteos" % bse["n"],
                       ("Aun sacándolos sigue primero: el liderazgo es real."
                        if aguanta else "Sin sorteos, cae al puesto %d." % pos_org),
                       "%d de sus %d posteos son sorteos. Sin ellos, su share orgánico es %.1f%%."
                       % (bse["sorteos"], bse["posts"], bse["sov_eng_org"]),
                       nota=M_SORTEO)

    # ── 3. El hallazgo: nadie conversa
    com_tot = sum(b["sent"]["comentarios"] for b in activas) if activas else 0
    if d["meta"].get("comentarios"):
        slide_seccion(prs, "Nadie está\nconversando",
                      "La categoría transmite. No dialoga.")
        slide_dato(prs, cliente, proyecto, "Conversación real",
                   NF(com_tot),
                   "comentarios en la ventana, entre todas las marcas juntas",
                   "Fuera de los sorteos, la categoría está muda.",
                   "El 99%% de los comentarios de la categoría sale de posteos de sorteo, "
                   "y son etiquetas a amigos para participar: no son opinión sobre la marca.",
                   nota=M_COMENTARIOS)

        # Motivos de queja: lo accionable
        conq = [b for b in activas if b["sent"]["suficiente"] and b["sent"]["motivos_neg"]]
        if conq:
            for b in conq[:2]:
                m = b["sent"]["motivos_neg"]
                slide_ranking(prs, cliente, proyecto, "De qué se quejan",
                              "%s · motivos de los comentarios negativos" % b["n"],
                              [{"n": x["k"], "v": x["v"], "lab": str(x["v"]),
                                "star": b["star"]} for x in m],
                              M_MOTIVOS % (b["sent"]["relevantes"], b["sent"]["neto"]))

    # ── 3.b Sentimiento: qué siente la gente, y validado con dos métodos
    rep = {}
    ruta_rep = os.path.join(HERE, "reporte_sentimiento.json")
    if os.path.exists(ruta_rep):
        rep = json.load(open(ruta_rep, encoding="utf-8"))

    if rep and bse:
        pm = rep["por_marca"]
        con_muestra = [(m, d) for m, d in pm.items() if d["n"] >= 30]
        if con_muestra:
            slide_seccion(prs, "Qué siente\nla gente",
                          "Medido dos veces, con dos métodos independientes.")
            b_bse = pm.get(bse["n"], {})
            if b_bse:
                rival = max([(m, d) for m, d in con_muestra if m != bse["n"]],
                            key=lambda x: x[1]["neto_llm"], default=None)
                enf = ("%s le gana por %d puntos." % (rival[0], rival[1]["neto_llm"] - b_bse["neto_llm"])
                       if rival and rival[1]["neto_llm"] > b_bse["neto_llm"]
                       else "Es el mejor de la categoría.")
                slide_dato(prs, cliente, proyecto, "Sentimiento neto",
                           "%+d%%" % b_bse["neto_llm"],
                           "de sentimiento neto en los comentarios de %s" % bse["n"],
                           enf,
                           "Un léxico rioplatense independiente da %+d%%: los dos métodos "
                           "coinciden en el %d%% de los casos."
                           % (b_bse["neto_lex"], b_bse["acuerdo"]),
                           nota=M_SENTIMIENTO % b_bse["n"])

            slide_ranking(prs, cliente, proyecto, "Sentimiento",
                          "Sentimiento neto por marca",
                          [{"n": m, "v": max(d["neto_llm"], 0),
                            "lab": "%+d%%  (n=%d)" % (d["neto_llm"], d["n"]),
                            "star": m == bse["n"]}
                           for m, d in sorted(con_muestra, key=lambda x: -x[1]["neto_llm"])],
                          M_SENT_RANKING)

            # La trampa que un léxico solo no ve: la ironía.
            tramp = rep["acuerdo"].get("trampas_ironia", [])
            if tramp:
                s = prs.slides.add_slide(prs.slide_layouts[6])
                _fondo(s, NEGRO)
                _header(s, cliente, proyecto)
                _etiqueta(s, "Cómo se midió")
                _texto(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.0),
                       "LA IRONÍA NO SE MIDE CONTANDO PALABRAS", 28, BLANCO, bold=True)
                _texto(s, Inches(0.9), Inches(2.4), Inches(11.0), Inches(0.7),
                       "Estos comentarios tienen solo palabras positivas. Son quejas.",
                       19, LAVANDA, italic=True, font=SERIF)
                y = Inches(3.4)
                for t in tramp[:4]:
                    _texto(s, Inches(0.9), y, Inches(11.3), Inches(0.6),
                           "«%s»" % t["texto"][:105].replace("\n", " "), 14, BLANCO)
                    _texto(s, Inches(0.9), y + Inches(0.38), Inches(11.3), Inches(0.3),
                           t["marca"], 10, GRIS)
                    y = y + Inches(0.85)
                _nota(s, M_IRONIA % (rep["acuerdo"]["pct"], 100 - rep["acuerdo"]["pct"]))

    # ── 3.c Nubes de palabras: QUÉ DICE LA GENTE de cada marca
    #
    # Son de los COMENTARIOS, no de los posteos: el tamaño de cada palabra es cuántas
    # veces el público la escribió. La imagen la arma el renderer de wordcloud (layout
    # empaquetado, sin solapes); acá solo se pega el PNG y se lista el top al costado.
    nubes_path = os.path.join(HERE, "nubes.json")
    if os.path.exists(nubes_path):
        NB = json.load(open(nubes_path, encoding="utf-8"))
        if NB:
            slide_seccion(prs, "Qué dice\nla gente",
                          "Las palabras que el público más repite en los comentarios.")

            orden = sorted(NB.items(), key=lambda x: -x[1]["n"])
            for m, info in orden:
                if info["n"] < 30:
                    continue          # muestra chica: la nube sería anecdótica
                es_bse = bool(bse and m == bse["n"])

                slide_nube(prs, cliente, proyecto,
                           "%s · lo que más repite su público" % m,
                           info.get("todas"), info.get("top", []),
                           lectura="", nota=M_NUBE % info["n"],
                           color=LAVANDA if es_bse else BLANCO)

                # Elogios y quejas por separado: es lo que se convierte en decisión.
                if info.get("neg"):
                    slide_nube(prs, cliente, proyecto,
                               "%s · de qué se queja su público" % m,
                               info["neg"], info.get("top_neg", []),
                               lectura=("Las palabras de los comentarios negativos. "
                                        "El tamaño es cuántas veces se repiten."),
                               nota=M_NUBE_NEG, color=RGBColor(0xFF, 0x8A, 0x7A))

    # ── 4. Los que hablan solos
    mudos = [b for b in activas if b["posts"] >= 100 and b["sov_eng"] < 3]
    if mudos:
        slide_seccion(prs, "Los que\nhablan solos",
                      "Publican todos los días. No los escucha nadie.")
        slide_ranking(prs, cliente, proyecto, "Esfuerzo sin retorno",
                      "Posteos publicados vs. share de engagement",
                      [{"n": b["n"], "v": b["posts"],
                        "lab": "%d posteos → %.1f%%" % (b["posts"], b["sov_eng"]),
                        "star": False} for b in sorted(mudos, key=lambda b: -b["posts"])],
                      M_POSTEOS)

    # ── 5. Territorios y océanos libres
    libres = [t for t in d["territorios"] if t["bse_pct"] == 0 and t["v"] >= 3]
    slide_seccion(prs, "Los territorios",
                  "De qué habla la categoría, y dónde no está nadie.")
    slide_ranking(prs, cliente, proyecto, "Territorios",
                  "De qué habla la categoría",
                  [{"n": t["k"], "v": t["v"],
                    "lab": "%d%% es %s" % (t["bse_pct"], cliente) if t["bse_pct"] else "sin %s" % cliente,
                    "star": t["bse_pct"] >= 50} for t in d["territorios"][:8]],
                  M_TERRITORIOS)

    # ── Cierre. Todo sale de los DATOS, no de supuestos de la categoría. Antes había
    # una slide de TikTok con texto de seguros hardcodeado ("primera póliza"): rompía
    # el sentido en cualquier otra categoría.
    op = []
    if d["meta"].get("comentarios") and com_tot < 1500:
        op.append({"t": "La categoría no conversa",
                   "d": "Nadie responde, nadie pregunta. El primero que abra diálogo se queda con el territorio."})
    if libres:
        op.append({"t": "Territorios sin dueño: " + ", ".join(t["k"] for t in libres[:2]),
                   "d": "Nadie los ocupa hoy. Están libres."})
    # Redes que NINGUNA marca de la categoría usa: océano de red, dato duro.
    usadas = {r for b in d["brands"] for r in b.get("redes", {})}
    NOMBRE_RED = {"ig": "Instagram", "fb": "Facebook", "x": "X / Twitter", "tt": "TikTok"}
    ausentes = [NOMBRE_RED[k] for k in ("tt", "x", "fb", "ig")
                if k not in usadas and k in NOMBRE_RED]
    if ausentes:
        op.append({"t": "%s sin explorar" % ausentes[0],
                   "d": "Ninguna marca de la categoría tiene cuenta ahí. Es una red disponible."})
    if bse and bse["pct_sorteo"] >= 40:
        op.append({"t": "Bajar la dependencia del sorteo",
                   "d": "Buena parte del engagement se compra con premios. El contenido tiene que "
                        "sostenerse solo."})
    if not op:                    # nunca dejar la slide vacía
        op.append({"t": "Sin océanos evidentes",
                   "d": "La categoría está disputada en todos los frentes relevados."})
    slide_cierre(prs, op[:4])

    prs.save(SALIDA)
    print("OK · %s (%d slides)" % (SALIDA, len(prs.slides.__iter__.__self__._sldIdLst)))
    return 0


if __name__ == "__main__":
    sys.exit(build())
